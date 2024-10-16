import os
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import shutil

import chromadb
from chromadb import Client

import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from llama_index.core import StorageContext
from llama_index.core import VectorStoreIndex, node_parser
from llama_index.core import  ( Document, GPTVectorStoreIndex, ServiceContext, SimpleDirectoryReader )
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.core import Settings
from ui import css, user_template, bot_template
# from llama_index.llms import OpenAI
# from llama_index.core import SimpleTextSplitter



# Function to process PDFs and create a vector store
def get_vectorstore(pdf_docs):
    text_chunks = []
    for pdf in pdf_docs:
        reader = PdfReader(pdf)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        text_chunks.append(text)

    documents = [Document(text=chunk) for chunk in text_chunks]
    embeddings = OpenAIEmbedding()
    vectorstore = VectorStoreIndex.from_documents(documents=documents, embedding=embeddings)
    return vectorstore



# For dealing with uploaded documents [ pdf, csv,excl,pptx,docx,]
def get_vectorstore_uploaded_docs(uploaded_docs):
    text_chunks = []

    for doc in uploaded_docs:
        file_extension = os.path.splitext(doc.name)[1].lower()

        if file_extension == '.pdf':
            # Handle PDF files
            reader = PdfReader(doc)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            text_chunks.append(text)

        elif file_extension == '.docx':
            # Handle DOCX files
            docx = DocxDocument(doc)
            text = ""
            for paragraph in docx.paragraphs:
                text += paragraph.text + "\n"
            text_chunks.append(text)

        elif file_extension == '.pptx':
            # Handle PPTX files
            presentation = Presentation(doc)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            text_chunks.append(text)

        elif file_extension in ['.csv', '.xls', '.xlsx']:
            # Handle CSV and Excel files
            if file_extension == '.csv':
                df = pd.read_csv(doc)
            else:
                df = pd.read_excel(doc)

            # Convert DataFrame to text (you may want to adjust this format)
            text = df.to_string(index=False)
            text_chunks.append(text)

        else:
            print(f"Unsupported file type: {file_extension}")

    # Create documents and embeddings
    documents = [Document(text=chunk) for chunk in text_chunks]
    embeddings = OpenAIEmbedding()
    vectorstore = VectorStoreIndex.from_documents(documents=documents, embedding=embeddings)
    
    return vectorstore



#  For getting docs direct from data directory 
def get_vectorstore_all_docs():
    embeddings = OpenAIEmbedding()
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    # documents = [Document(text=chunk) for chunk in text_chunks]
    # documents = SimpleDirectoryReader(text_chunks).load_data()
    documents = SimpleDirectoryReader("/home/ansh/work/thum/Llamindex-Projects/poc/data/").load_data()

    vectorstore = VectorStoreIndex.from_documents(documents=documents, embedding=embeddings)
    return vectorstore





#  Saves uplaoded_docs  by users directly to data directory 
def save_uploaded_files(uploaded_files, directory="data"):
    # Check if the 'data' directory exists, if so, clear its contents
    if os.path.exists(directory):
        shutil.rmtree(directory)  # Delete the directory and its contents
    os.makedirs(directory)  # Recreate the directory

    # Save each uploaded file to the 'data' directory
    for file in uploaded_files:
        with open(os.path.join(directory, file.name), "wb") as f:
            f.write(file.getbuffer())




# Function to handle user input and generate chatbot response
def handle_userinput(user_question):
    # Check if the conversation has been initialized
    if st.session_state.conversation is None:
        st.warning("Please upload documents first before asking questions.")
        return

    # Get the bot response from the vector store query engine
    response = st.session_state.conversation.query(user_question)

    # Append user and bot responses to chat history
    st.session_state.chat_history.append(("user", user_question))
    st.session_state.chat_history.append(("bot", str(response)))  # Convert response to string

    # Display the conversation with user question on top and bot answer below
    for i in range(0, len(st.session_state.chat_history), 2):
        if i < len(st.session_state.chat_history):
            user_message = st.session_state.chat_history[i][1]  # User's question
            st.markdown(user_template.replace("{{MSG}}", user_message), unsafe_allow_html=True)
        if i + 1 < len(st.session_state.chat_history):
            bot_message = st.session_state.chat_history[i + 1][1]  # Bot's response
            st.markdown(bot_template.replace("{{MSG}}", bot_message), unsafe_allow_html=True)






# Main function for the Streamlit app
def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple Documents", page_icon=":books:")
    st.markdown(css, unsafe_allow_html=True)

    # Initialize conversation and chat history in session state
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    # Container for the chat
    st.markdown("<div class='container'><h2>Chat with your Documents</h2></div>", unsafe_allow_html=True)

    # Create a form to handle user input
    with st.form(key='chat_input_form', clear_on_submit=True):
        user_input = st.text_input("Ask a question about your documents:")
        submit_button = st.form_submit_button("Send")

    # Clear button using HTML
    if st.button("Clear Chat", key='clear_chat', help="Clear chat history"):
        # Clear the chat history and reset conversation state
        st.session_state.chat_history = []
        st.session_state.conversation = None
        st.success("Chat history cleared!")

    if submit_button and user_input:
        handle_userinput(user_input)

    with st.sidebar:
        st.subheader("Your documents")
        # pdf_docs = st.file_uploader("Upload your PDFs here and click on 'Process'", accept_multiple_files=True)
        uploaded_docs = st.file_uploader("Upload your Docs. here and click on 'Process'", accept_multiple_files=True)
        
          
        if uploaded_docs:
            if st.button("Process"):
                with st.spinner("Processing..."):
                    # Clear chat history when new documents are processed
                    st.session_state.chat_history = []

                    # Save the uploaded files
                    save_uploaded_files(uploaded_docs)

                    # Create the vector store from the uploaded PDFs
                    # vectorstore = get_vectorstore(pdf_docs)
                    # vectorstore =  get_vectorstore_uploaded_docs(uploaded_docs) # use when you have pdf,excl,csv,etc  docs.
                    vectorstore = get_vectorstore_all_docs()     # use when you have various different docs.

                    # Initialize the conversation (query engine)
                    st.session_state.conversation = vectorstore.as_query_engine()

        else:
            st.warning("Please upload files to proceed.")



if __name__ == '__main__':
    main()

